import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- CONFIGURAÇÃO ---
PASTA_RESULTADOS = "resultados_analise_IJUI_FINAL"
NOME_ARQUIVO_PPT = "Apresentacao_Saude_Ijui.pptx"

print(f"--- GERANDO APRESENTAÇÃO AUTOMÁTICA ---")

# Verifica se a pasta existe
if not os.path.exists(PASTA_RESULTADOS):
    print(f"❌ ERRO: A pasta '{PASTA_RESULTADOS}' não foi encontrada.")
    exit()

# Cria a apresentação
prs = Presentation()

# --- FUNÇÕES AUXILIARES ---

def slide_titulo(titulo, subtitulo):
    """Cria um slide de capa."""
    slide_layout = prs.slide_layouts[0] # 0 é o layout de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titulo
    subtitle.text = subtitulo

def slide_texto_imagem(titulo, texto_lista, nome_imagem=None):
    """Cria um slide com título, bullets e uma imagem opcional."""
    slide_layout = prs.slide_layouts[1] # 1 é Título + Conteúdo
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    title = slide.shapes.title
    title.text = titulo
    
    # Texto (Bullets)
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    for item in texto_lista:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # Imagem (se houver)
    if nome_imagem:
        caminho_img = os.path.join(PASTA_RESULTADOS, nome_imagem)
        if os.path.exists(caminho_img):
            # Adiciona imagem no canto direito inferior (ajuste conforme necessário)
            left = Inches(4.5)
            top = Inches(2.5)
            height = Inches(4.0)
            slide.shapes.add_picture(caminho_img, left, top, height=height)

def slide_tabela(titulo, nome_csv, top_n=10):
    """Cria um slide com uma tabela baseada em um CSV."""
    caminho_csv = os.path.join(PASTA_RESULTADOS, nome_csv)
    if not os.path.exists(caminho_csv):
        print(f"Aviso: {nome_csv} não encontrado.")
        return

    # Lê o CSV
    df = pd.read_csv(caminho_csv, sep=';')
    df = df.head(top_n) # Pega apenas as primeiras linhas para caber no slide

    slide_layout = prs.slide_layouts[5] # 5 é Título Apenas
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titulo

    # Configurações da tabela
    rows, cols = df.shape
    rows += 1 # +1 para o cabeçalho
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Preenche cabeçalho
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102) # Azul escuro
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255) # Branco
        paragraph.font.bold = True

    # Preenche dados
    for i, row in df.iterrows():
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            # Formata valores numéricos grandes
            if isinstance(val, (int, float)) and val > 1000:
                cell.text = f"{val:,.0f}".replace(",", ".")
            else:
                cell.text = str(val)
            
            cell.text_frame.paragraphs[0].font.size = Pt(10) # Fonte menor para caber

# --- MONTAGEM DA APRESENTAÇÃO ---

# 1. Capa
slide_titulo("Relatório de Gestão em Saúde", "Análise de Dados SIASUS - Município de Ijuí/RS")

# 2. Volume e Tendência
slide_texto_imagem(
    "1. Volume e Tendência Temporal",
    [
        "Análise da evolução da demanda ambulatorial.",
        "Destaque para a concentração de dados em meses específicos.",
        "Visão geral do volume processado."
    ],
    "grafico_tendencia_demanda.png"
)
slide_tabela("Detalhe do Volume (Por Trimestre)", "analise_3_1_trimestre.csv")

# 3. Estabelecimentos
slide_tabela("2. Top 10 Estabelecimentos (Produção)", "analise_3_2_estabelecimentos.csv", top_n=10)

# 4. Perfil Demográfico (Idade)
slide_texto_imagem(
    "3. Perfil Etário dos Pacientes",
    [
        "Distribuição dos atendimentos por faixa etária.",
        "Identificação dos grupos de maior demanda (ex: Idosos)."
    ],
    "grafico_faixa_etaria.png"
)
slide_tabela("Dados Demográficos", "analise_3_3_demografico_sexo_idade.csv")

# 5. Origem dos Pacientes
slide_texto_imagem(
    "4. Fluxo Regional: De onde vêm os pacientes?",
    [
        "Proporção entre munícipes de Ijuí e pacientes de fora.",
        "Papel de Ijuí como polo regional de saúde."
    ],
    "grafico_origem_pacientes.png"
)
slide_tabela("Top Municípios de Origem", "analise_3_4_origem_pacientes.csv", top_n=8)

# 6. Financeiro
slide_tabela("5. Resumo Financeiro (Aprovado vs Produzido)", "analise_3_5_financeiro.csv")

# 7. Áreas Críticas (Oncologia, Saúde Mental, Diabetes)
slide_tabela("6. Foco: Oncologia (Top Procedimentos)", "analise_3_6_oncologia.csv", top_n=6)
slide_tabela("6. Foco: Saúde Mental", "analise_3_6_saude_mental.csv", top_n=6)
slide_tabela("6. Foco: Diabetes", "analise_3_6_diabetes.csv", top_n=6)

# 8. Tendências em Idosos
slide_texto_imagem(
    "7. Tendências: Idosos (Cardio vs. Onco)",
    [
        "Acompanhamento de áreas críticas na população 60+.",
        "Comparativo entre Cardiologia e Oncologia."
    ],
    "grafico_tendencia_idosos.png"
)

# 9. Encerramento
slide_titulo("Conclusão", "Dados extraídos da base SIASUS (MySQL)")

# --- SALVAR ---
prs.save(NOME_ARQUIVO_PPT)
print(f"\n✅ Sucesso! Apresentação salva como: {NOME_ARQUIVO_PPT}")